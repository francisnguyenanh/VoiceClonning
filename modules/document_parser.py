"""
document_parser.py — Extract plain text from supported document formats.

Supported: .txt, .docx, .pptx, .xlsx, .pdf
"""
from __future__ import annotations

import logging
from pathlib import Path

logger = logging.getLogger(__name__)


def parse_document(filepath: str | Path) -> str:
    """
    Parse a document and return its text content.

    Args:
        filepath: Path to the document file.

    Returns:
        Plain text string extracted from the document.

    Raises:
        ValueError: If the file extension is not supported.
        FileNotFoundError: If the file does not exist.
    """
    path = Path(filepath)
    if not path.exists():
        raise FileNotFoundError(f"File not found: {path}")

    ext = path.suffix.lower()
    parsers = {
        ".txt":  _parse_txt,
        ".docx": _parse_docx,
        ".pptx": _parse_pptx,
        ".xlsx": _parse_xlsx,
        ".pdf":  _parse_pdf,
    }

    if ext not in parsers:
        raise ValueError(f"Unsupported file type: {ext}. Supported: {', '.join(parsers)}")

    logger.info("Parsing document: %s (type=%s)", path.name, ext)
    return parsers[ext](path)


# ---------------------------------------------------------------------------
# Format-specific parsers
# ---------------------------------------------------------------------------

def _parse_txt(path: Path) -> str:
    for encoding in ("utf-8", "utf-8-sig", "utf-16", "cp1252"):
        try:
            return path.read_text(encoding=encoding)
        except (UnicodeDecodeError, LookupError):
            continue
    raise ValueError(f"Cannot decode text file: {path}")


def _parse_docx(path: Path) -> str:
    from docx import Document  # python-docx
    doc = Document(str(path))
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    return "\n".join(paragraphs)


def _parse_pptx(path: Path) -> str:
    from pptx import Presentation  # python-pptx
    prs = Presentation(str(path))
    lines: list[str] = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        lines.append(f"[Slide {slide_num}]")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
    return "\n".join(lines)


def _parse_xlsx(path: Path) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    lines: list[str] = []
    for sheet in wb.worksheets:
        lines.append(f"[Sheet: {sheet.title}]")
        for row in sheet.iter_rows(values_only=True):
            row_text = ", ".join(str(cell) for cell in row if cell is not None)
            if row_text.strip():
                lines.append(row_text)
    wb.close()
    return "\n".join(lines)


def _parse_pdf(path: Path) -> str:
    from pdfminer.high_level import extract_text
    text = extract_text(str(path))
    return text.strip() if text else ""
